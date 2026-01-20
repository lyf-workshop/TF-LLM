import json
import logging
from typing import Any

import yaml
from ppt_template_model import PageConfig
from pptx import Presentation
from utils import delete_slide, delete_slide_range, duplicate_slide, move_slide


def fill_template_with_yaml_config(template_path, output_path, json_data, yaml_config: dict[str, Any]):
    page_config = PageConfig(yaml_config)
    prs = Presentation(template_path)
    data = json.loads(json_data)
    slides_data = data.get("slides", [])

    if not isinstance(slides_data, list):
        raise ValueError("JSON data must contain a 'slides' list")

    for slide_data in slides_data:
        slide_type = slide_data.get("type")
        if not slide_type:
            logging.warning("Skipped slide without type definition: %s", slide_data)
            continue

        template_index = page_config.type_map.get(slide_type)
        if template_index is None or template_index >= len(prs.slides):
            logging.warning("No template found for slide type '%s'", slide_type)
            continue

        template_slide = prs.slides[template_index]
        if slide_type in ("title", "acknowledgement"):
            target_slide = template_slide
        else:
            target_slide = duplicate_slide(prs, template_slide)

        page_config.render(target_slide, slide_data)

    delete_slide_range(prs, range(2, 12))
    delete_slide(prs, 0)
    move_slide(prs, 1, len(prs.slides) - 1)
    prs.save(output_path)


def extract_json(content):
    """
    Extract the json data from the given content.
    """
    # extract content within "```json" and "```"
    json_data = content.split("```json")[1].split("```")[0]
    return json_data


if __name__ == "__main__":
    import argparse
    import datetime

    parser = argparse.ArgumentParser()
    parser.add_argument("-t", "--template", type=str, default="templates/0.pptx")
    default_output_filename = f"output-{datetime.datetime.now().strftime('%Y-%m-%d-%H-%M-%S')}.pptx"
    parser.add_argument("-o", "--output", type=str, default=default_output_filename)
    parser.add_argument("-i", "--input", type=str, required=True)
    parser.add_argument("--cache_dir", type=str, default=".temp")
    parser.add_argument("--yaml_config", type=str, default="yaml_example.yaml")
    args = parser.parse_args()

    # set env var UTU_PPT_CACHE_DIR
    import os

    os.environ["UTU_PPT_CACHE_DIR"] = args.cache_dir

    logging.basicConfig(level=logging.INFO)
    template = args.template
    output = args.output
    input_json = args.input
    with open(input_json) as f:
        content = f.read()
    json_data = extract_json(content)

    with open(args.yaml_config) as f:
        yaml_config = yaml.safe_load(f)

    fill_template_with_yaml_config(template, output, json_data, yaml_config)
